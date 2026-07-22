#!/usr/bin/env python3
"""Build the AI-era rebuild of the Real-Time Systems lecture deck.

Template-free (default python-pptx presentation, 16:9). Concept content re-authored
with a verification-vs-validation lens; teaching automata redrawn natively; genuine
screenshots reused from assets/. See AI-rebuild-design.md for the full design.
"""
import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")

# ---------- design system ----------
NAVY   = RGBColor(0x1E, 0x27, 0x61)
INK    = RGBColor(0x20, 0x24, 0x33)
TEAL   = RGBColor(0x02, 0x80, 0x90)
AMBER  = RGBColor(0xC9, 0x7A, 0x12)      # AI-era callout accent
AMBERBG= RGBColor(0xFB, 0xF1, 0xDE)
LIGHT  = RGBColor(0xF7, 0xF8, 0xFB)
CARD   = RGBColor(0xEC, 0xF0, 0xF6)
MUTED  = RGBColor(0x55, 0x5E, 0x6E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LOCFILL= RGBColor(0xE7, 0xEE, 0xF8)
DARKBG = RGBColor(0x16, 0x1C, 0x3A)

HEAD = "Cambria"
BODY = "Calibri"
MONO = "Consolas"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ---------- low-level helpers ----------
def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = DARKBG if dark else LIGHT
    return s

def _set(run, size, color, bold=False, italic=False, font=BODY):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.name = font; run.font.color.rgb = color

def box(s, x, y, w, h, anchor=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tb, tf

def para(tf, text, size=15, color=INK, bold=False, italic=False, font=BODY,
         align=PP_ALIGN.LEFT, before=0, after=6, first=False, line=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(before); p.space_after = Pt(after)
    p.line_spacing = line
    r = p.add_run(); r.text = text
    _set(r, size, color, bold, italic, font)
    return p

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, rounded=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def shape_text(shp, lines, size=12, color=INK, bold=False, font=BODY, align=PP_ALIGN.CENTER):
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = 0; p.space_after = 0; p.line_spacing = 1.0
        r = p.add_run(); r.text = ln
        _set(r, size, color, bold, font=font)

# ---------- slide furniture ----------
def title(s, text, dark=False, kicker=None):
    col = WHITE if dark else NAVY
    y = 0.5
    if kicker:
        _, tf = box(s, 0.7, 0.42, 12, 0.32)
        para(tf, kicker.upper(), size=12.5, color=TEAL, bold=True, first=True, after=0)
        y = 0.78
    # small square motif
    rect(s, 0.7, y + 0.11, 0.16, 0.16, fill=TEAL)
    _, tf = box(s, 0.98, y, 11.9, 1.0)
    para(tf, text, size=29, color=col, bold=True, font=HEAD, first=True, after=0, line=1.0)

def pagenum(s, n, dark=False):
    _, tf = box(s, 12.2, 7.06, 1.0, 0.3)
    para(tf, str(n), size=10, color=(RGBColor(0x9F, 0xA8, 0xC4) if dark else MUTED),
         align=PP_ALIGN.RIGHT, first=True, after=0)

def _bullets_height(items, w, base, gap):
    """Conservative estimate of rendered bullet-block height, in inches."""
    total = 0.0
    for i, (lvl, txt) in enumerate(items):
        if txt.startswith("`") and txt.endswith("`"): txt = txt[1:-1]
        if lvl == 0:   fs, pre = base, 3
        elif lvl == 1: fs, pre = base - 2.5, 8
        else:          fs, pre = base - 3, 0
        cpl = max(6.0, w * 144.0 / fs)              # chars per line (~0.5*fs/72 wide)
        lines = max(1, math.ceil((len(txt) + pre) / cpl))
        total += lines * fs * 1.2                   # line height (pt)
        total += (0 if i == 0 else (8 if lvl == 0 else 1))   # space before
        total += gap if lvl == 0 else gap - 1                # space after
    return total / 72.0

def bullets(s, x, y, w, items, h=5.0, gap=6, base=18, bottom=6.7,
            grow=True, max_base=26, min_base=15):
    """items: list of (level, text). Wrap `..` in backticks for mono.
    When grow=True, pick the largest font that fills (bottom - y) and spread
    any leftover space across the gaps for an even fill."""
    if grow:
        avail = max(1.0, bottom - y)
        base = min_base
        b = max_base
        while b >= min_base:
            if _bullets_height(items, w, b, max(6.0, b * 0.5)) <= avail * 0.97:
                base = b; break
            b -= 0.5
        gap = max(6.0, base * 0.5)
        used = _bullets_height(items, w, base, gap)
        if len(items) > 1 and used < avail:                 # distribute slack into gaps
            gap = min(base * 1.7, gap + (avail - used) * 72.0 / (len(items) - 1) * 0.7)
    _, tf = box(s, x, y, w, max(h, bottom - y))
    for i, (lvl, txt) in enumerate(items):
        mono = txt.startswith("`") and txt.endswith("`")
        if mono: txt = txt[1:-1]
        if lvl == 0:
            r = para(tf, "", size=base, color=INK, bold=True, first=(i == 0),
                     before=(0 if i == 0 else 8), after=gap)
            run = r.runs[0]; run.text = "▸  "; _set(run, base, TEAL, True)
            run2 = r.add_run(); run2.text = txt
            _set(run2, base, INK, True, font=(MONO if mono else BODY))
        elif lvl == 1:
            r = para(tf, "", size=base-2.5, color=MUTED, first=(i == 0),
                     before=1, after=gap-1)
            run = r.runs[0]; run.text = "      –  "; _set(run, base-2.5, MUTED)
            run2 = r.add_run(); run2.text = txt
            _set(run2, base-2.5, MUTED, font=(MONO if mono else BODY))
        else:  # lvl 2 plain / note
            para(tf, txt, size=base-3, color=MUTED, italic=True, first=(i == 0),
                 before=1, after=gap-1)
    return tf

def callout(s, x, y, w, h, head, body_lines, accent=AMBER, bg=AMBERBG):
    c = rect(s, x, y, w, h, fill=bg, line=accent, lw=1.25, rounded=True)
    rect(s, x, y, 0.11, h, fill=accent)  # left accent strip
    _, tf = box(s, x + 0.28, y + 0.16, w - 0.5, h - 0.3)
    para(tf, head, size=12.5, color=accent, bold=True, first=True, after=4)
    for ln in body_lines:
        para(tf, ln, size=13.5, color=INK, after=3, line=1.06)
    return c

def add_image(s, path, x, y, w=None, h=None, frame=True):
    im = Image.open(path); iw, ih = im.size; ar = iw / ih
    if w and not h: h = w / ar
    if h and not w: w = h * ar
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if frame:
        pic.line.color.rgb = RGBColor(0xCF, 0xD6, 0xE2); pic.line.width = Pt(1)
    return pic

# ---------- automaton primitives ----------
# label kinds → colour: guard (green), action/sync (navy), update/assignment (amber-brown)
GUARD_C = RGBColor(0x15, 0x86, 0x5A)
ACT_C   = NAVY
ASGN_C  = RGBColor(0xB0, 0x55, 0x0A)
LBL = {'g': GUARD_C, 'a': ACT_C, 'u': ASGN_C}

def _arrow(shp, tail=True, head=False):
    ln = shp.line._get_or_add_ln()
    if head:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if tail:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))

def _norm(labels):
    out = []
    for it in labels or []:
        if isinstance(it, (tuple, list)):
            out.append((it[0], it[1]))
        else:
            out.append(('a', str(it)))
    return out

def _place_labels(s, cx, top_y, items, w=3.0):
    """Render each (kind,text) on its own centred, colour-coded monospace line."""
    items = _norm(items)
    _, tf = box(s, cx - w / 2, top_y, w, 0.26 * len(items) + 0.12)
    for i, (k, t) in enumerate(items):
        para(tf, t, size=13, color=LBL.get(k, ACT_C), bold=(k == 'a'),
             font=MONO, align=PP_ALIGN.CENTER, first=(i == 0), after=0, line=1.05)

def legend(s, x=0.98, y=6.98):
    _, tf = box(s, x, y, 6.0, 0.3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "labels:  "; _set(r, 12, MUTED)
    for k, name, sep in [('g', 'guard', '    '), ('a', 'action', '    '), ('u', 'update', '')]:
        r = p.add_run(); r.text = name + sep
        _set(r, 12, LBL[k], bold=(k == 'a'), font=MONO)

def loc(s, cx, cy, name, r=0.42, initial=False, tag=None, tag_side='above', fill=LOCFILL):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
                           Inches(2 * r), Inches(2 * r))
    o.fill.solid(); o.fill.fore_color.rgb = fill
    o.line.color.rgb = NAVY; o.line.width = Pt(1.75); o.shadow.inherit = False
    shape_text(o, [name], size=14, color=NAVY, bold=True)
    if tag:
        if tag_side == 'right':
            _, tf = box(s, cx + r + 0.08, cy - 0.13, 1.1, 0.26)
            para(tf, tag, size=11, color=TEAL, bold=True, italic=True, first=True, after=0)
        elif tag_side == 'below':
            _, tf = box(s, cx - r, cy + r + 0.04, 2 * r, 0.26)
            para(tf, tag, size=11, color=TEAL, bold=True, italic=True,
                 align=PP_ALIGN.CENTER, first=True, after=0)
        else:
            _, tf = box(s, cx - r, cy - r - 0.28, 2 * r, 0.26)
            para(tf, tag, size=11, color=TEAL, bold=True, italic=True,
                 align=PP_ALIGN.CENTER, first=True, after=0)
    if initial:
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(cx - r - 0.55), Inches(cy),
                                   Inches(cx - r), Inches(cy))
        c.line.color.rgb = NAVY; c.line.width = Pt(1.75); _arrow(c)
    return (cx, cy, r, o)

def _unit(ax, ay, bx, by):
    dx, dy = bx - ax, by - ay
    d = (dx * dx + dy * dy) ** 0.5 or 1
    return dx / d, dy / d, d

def edge(s, a, b, labels=None, off=0.0, side='above', lab_shift=0.0):
    ax, ay, ar = a[0], a[1], a[2]; bx, by, br = b[0], b[1], b[2]
    ux, uy, _ = _unit(ax, ay, bx, by)
    px, py = -uy, ux
    a_along = max(0.05, ar * ar - off * off) ** 0.5
    b_along = max(0.05, br * br - off * off) ** 0.5
    sx, sy = ax + px * off + ux * a_along, ay + py * off + uy * a_along
    ex, ey = bx + px * off - ux * b_along, by + py * off - uy * b_along
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(sx), Inches(sy),
                               Inches(ex), Inches(ey))
    c.line.color.rgb = NAVY; c.line.width = Pt(1.5); c.shadow.inherit = False; _arrow(c)
    if labels:
        items = _norm(labels); n = len(items)
        mx = (sx + ex) / 2 + lab_shift; my = (sy + ey) / 2
        top = my - 0.10 - 0.26 * n if side == 'above' else my + 0.10
        _place_labels(s, mx, top, items)
    return c

def _freeform(s, pts, color=NAVY, width=1.5, arrow=True):
    ff = s.shapes.build_freeform(Inches(pts[0][0]), Inches(pts[0][1]), scale=1.0)
    ff.add_line_segments([(Inches(x), Inches(y)) for x, y in pts[1:]], close=False)
    shp = ff.convert_to_shape()
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(width)
    shp.shadow.inherit = False
    if arrow:
        _arrow(shp)
    return shp

def self_loop_curved(s, node, labels=None):
    """A curved self-loop that leaves the top of the node and returns to it."""
    cx, cy, r = node[0], node[1], node[2]
    R = 0.40
    Cx, Cy = cx, cy - r - 0.16
    pts = [(cx - 0.15, cy - r + 0.02)]                     # start on node boundary (top-left)
    for i in range(1, 20):                                  # arc over the top, gap at the bottom
        t = math.radians(150 + (240) * i / 20)              # 150° → 390°
        pts.append((Cx + R * math.cos(t), Cy + R * math.sin(t)))
    pts.append((cx + 0.15, cy - r + 0.02))                  # end on node boundary (top-right)
    _freeform(s, pts)
    if labels:
        top = (Cy - R) - 0.10 - 0.26 * len(_norm(labels))
        _place_labels(s, cx, top, labels)

def elbow_edge(s, a, b, labels=None, drop=0.9, end_dx=0.0, lab_side='above'):
    """Right-angle (elbow) edge routed below both nodes, arrow into b's bottom."""
    ax, ay, ar = a[0], a[1], a[2]; bx, by, br = b[0], b[1], b[2]
    y_level = max(ay + ar, by + br) + drop
    tx = bx + end_dx
    pts = [(ax, ay + ar), (ax, y_level), (tx, y_level), (tx, by + br + 0.01)]
    _freeform(s, pts)
    if labels:
        mx = (ax + tx) / 2
        n = len(_norm(labels))
        top = y_level - 0.10 - 0.26 * n if lab_side == 'above' else y_level + 0.08
        _place_labels(s, mx, top, labels)

def diagram_caption(s, x, y, w, text):
    _, tf = box(s, x, y, w, 0.4)
    para(tf, text, size=13.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
         first=True, after=0)

# =======================================================================
# SLIDES
# =======================================================================

# 1 — Title (dark)
s = slide(dark=True)
rect(s, 0, 0, 0.28, SH, fill=TEAL)
_, tf = box(s, 1.1, 2.35, 11, 2.2)
para(tf, "Formal Modeling and Analysis of", size=36, color=WHITE, bold=True, font=HEAD, first=True, after=0, line=1.02)
para(tf, "Real-Time Systems", size=36, color=WHITE, bold=True, font=HEAD, after=6, line=1.02)
para(tf, "…in the age of AI", size=20, color=RGBColor(0x8F,0xB9,0xC7), italic=True, font=HEAD, after=0)
_, tf = box(s, 1.12, 5.35, 11, 1.2)
para(tf, "Rong Gu", size=17, color=WHITE, bold=True, first=True, after=1)
para(tf, "ronggufly@gmail.com", size=14, color=RGBColor(0x9F,0xA8,0xC4), after=1)
para(tf, "Dec 16, 2025", size=13, color=RGBColor(0x9F,0xA8,0xC4), after=0)

# 2 — Who am I
s = slide()
title(s, "Who am I?")
bullets(s, 0.98, 1.75, 11.8, [
    (0, "Rong Gu — Formal Methods Developer at Prover (2025 – now)"),
    (1, "2010 / 2013: BSc & MSc in Software Engineering, Xi'an Jiaotong University, China"),
    (1, "2013–2017: Embedded software engineer, China"),
    (1, "2022: PhD in Computer Science, Mälardalen University (Seceleanu, Enoiu, Lundqvist)"),
    (1, "2022–2025: Postdoc, Mälardalen University (IDT) — formal modelling & analysis of embedded systems"),
    (0, "Contact:  ronggufly@gmail.com"),
], gap=16, base=21)
pagenum(s, 2)

# 3 — What is a model?
s = slide()
title(s, "What is a model?", kicker="Foundations")
bullets(s, 0.98, 1.9, 7.2, [
    (0, "Model — an abstract representation of reality"),
    (1, "leaves out detail to make the workings understandable"),
    (0, "A system model helps you…"),
    (1, "understand the functionality of the system"),
    (1, "communicate with customers and colleagues"),
    (1, "…and, if given precise semantics, reason about it rigorously"),
], gap=14, base=20)
callout(s, 8.35, 1.9, 4.35, 3.5, "AI is a model too — but a different kind",
        ["A large language model is itself a model of reality — but a statistical, non-semantic one you cannot feed to a prover for a guarantee.",
         "This course is about models you can reason about, and about checking the ones an AI builds for you."])
pagenum(s, 3)

# 4 — Perspectives & examples
s = slide()
title(s, "Models: perspectives & examples", kicker="Foundations")
bullets(s, 0.98, 1.55, 6.6, [
    (0, "Three perspectives on a system"),
    (1, "External — its context / environment"),
    (1, "Behavioural — what the system does over time"),
    (1, "Structural — its architecture"),
    (0, "Examples"),
    (1, "Data-flow models — how data is processed as it moves"),
    (1, "State-machine models — response to events"),
])
callout(s, 8.0, 1.7, 4.7, 2.5, "Is it a model?  Is an AI a model?",
        ["Both describe a system's behaviour.",
         "The difference is verifiability: a state machine has",
         "semantics you can check; an AI's internal 'model'",
         "does not. Keep that distinction in mind all lecture."],
        accent=TEAL, bg=CARD)
pagenum(s, 4)

# 5 — Formal models
s = slide()
title(s, "Formal models", kicker="Foundations")
bullets(s, 0.98, 1.7, 7.0, [
    (0, "A formal model has mathematical semantics"),
    (1, "every symbol and transition has a precise, unambiguous meaning"),
    (0, "That precision is what makes rigorous analysis possible"),
    (1, "you can prove properties, not just inspect diagrams"),
], gap=10)
callout(s, 8.2, 1.9, 4.5, 3.0, "The trust anchor",
        ["Formal semantics is exactly what lets you check",
         "someone else's model — including one produced by AI.",
         "",
         "No semantics → no proof → no basis for trust.",
         "This is why we still model formally when AI can",
         "generate artifacts in seconds."])
pagenum(s, 5)

# 6 — Model analysis: V&V (ANCHOR)
s = slide()
title(s, "Model analysis: Verification vs Validation", kicker="The spine of this course")
bullets(s, 0.98, 1.55, 5.9, [
    (0, "Analysis needs two inputs"),
    (1, "a system model, and a system requirement"),
    (0, "Verification — “did we build the thing right?”"),
    (1, "formal verification, simulation, testing"),
    (0, "Validation — “did we build the right thing?”"),
    (1, "user validation, test scenarios, judgement"),
], gap=7)
callout(s, 7.15, 1.55, 5.55, 4.2, "Where AI fits",
        ["VERIFICATION is symbolic and automatable —",
         "AI helps most here: draft models, TCTL properties,",
         "run tools.",
         "",
         "VALIDATION connects the model to real-world intent.",
         "It lives OUTSIDE the formal system, so AI cannot own",
         "it — the AI only has what you typed, not the intent.",
         "",
         "AI manipulates the symbols; it cannot supply the",
         "intent or the guarantee."])
pagenum(s, 6)

# 7 — Formal verification
s = slide()
title(s, "Formal verification", kicker="Verification")
bullets(s, 0.98, 1.6, 7.0, [
    (0, "Check ALL system behaviours — exhaustive, not sampled"),
    (0, "The property to check is stated in a logic"),
    (0, "Algorithmic verification (model checking)"),
    (1, "explore the state space, check it against logic properties"),
    (1, "scalability is still a genuine challenge"),
], gap=9)
callout(s, 8.2, 1.9, 4.5, 2.3, "Exhaustive ≠ probabilistic",
        ["An LLM gives a plausible answer, fast, with no guarantee. Model checking gives certainty over every reachable behaviour.",
         "Powerful ≠ certain."])
pagenum(s, 7)

# 8 — State-space explosion (reuse image)
s = slide()
title(s, "Model checking: the state-space explosion", kicker="Verification")
add_image(s, os.path.join(ASSETS, "state-explosion.png"), 0.98, 1.7, w=6.9)
bullets(s, 8.2, 1.7, 4.5, [
    (0, "The number of states grows explosively with components, clocks and variables"),
    (0, "This is an open research challenge…"),
    (1, "…orthogonal to how capable LLMs become"),
    (2, "AI does not dissolve state explosion — the reason model checking is still hard is not a reason AI can replace it."),
], gap=8, base=16.5)
pagenum(s, 8)

# 9 — NEW: Formal methods in the age of AI (thesis)
s = slide(dark=True)
title(s, "Formal methods in the age of AI", dark=True, kicker="Why this still matters")
# three columns
def thesis_card(x, head, lines):
    rect(s, x, 2.15, 3.7, 3.9, fill=RGBColor(0x22,0x2A,0x52), line=TEAL, lw=1.25, rounded=True)
    _, tf = box(s, x + 0.3, 2.4, 3.1, 3.5)
    para(tf, head, size=16, color=WHITE, bold=True, font=HEAD, first=True, after=8)
    for ln in lines:
        para(tf, ln, size=12.5, color=RGBColor(0xC7,0xD2,0xE8), after=4, line=1.08)
thesis_card(0.95, "AI lowers the cost",
    ["Building models and drafting", "properties — historically the", "barrier to formal methods —", "is now fast and assisted."])
thesis_card(4.82, "…and raises the stakes",
    ["If AI authors the artifact, you", "must verify it. The verifier is", "the trust anchor, and AI is a", "fallible author."])
thesis_card(8.68, "Your job shifts",
    ["From drawing automata to:", "specify precisely, validate", "against intent, and judge the", "counterexample. That is durable."])
pagenum(s, 9, dark=True)

# 10 — Deductive verification & theorem proving (merge 10+11)
s = slide()
title(s, "Deductive verification (theorem proving)", kicker="Verification")
bullets(s, 0.98, 1.6, 6.7, [
    (0, "Prove correctness from axioms and inference rules"),
    (1, "not fully automated — needs human guidance"),
    (0, "Classic syllogism"),
    (1, "All men are mortal;  Socrates is a man"),
    (1, "∴ Socrates is mortal"),
], gap=8)
callout(s, 7.9, 1.7, 4.8, 3.4, "Plausible reasoning vs sound proof",
        ["An LLM can imitate this reasoning — and often gets it",
         "right — but offers no guarantee of soundness.",
         "A proof does.",
         "",
         "Genuinely shifted by AI: LLMs + proof assistants are",
         "pushing autoformalization, making deductive",
         "verification more automatable than before."])
pagenum(s, 10)

# 11 — Model checking = A + F
s = slide()
title(s, "Model checking = model (A) + requirement (F)", kicker="Model checking")
bullets(s, 0.98, 1.7, 7.2, [
    (0, "A — the model: a network of timed automata"),
    (0, "F — the requirement: a temporal-logic formula, e.g."),
    (1, "Invariant — something bad never happens"),
    (1, "Reachability — something can eventually happen"),
    (1, "Liveness — something always eventually happens"),
    (2, "Tool: UPPAAL (Uppsala University & Aalborg University)"),
], gap=8)
callout(s, 8.4, 1.85, 4.3, 2.6, "The oracle",
        ["The model checker is the oracle of truth: given A and F it answers with certainty, and hands you a counterexample when F fails.",
         "It is where an AI-built model gets checked."], accent=TEAL, bg=CARD)
pagenum(s, 11)

# 12 — Real-time systems
s = slide()
title(s, "Real-time systems", kicker="Real-time systems")
bullets(s, 0.98, 1.55, 7.2, [
    (0, "Correctness depends on the logical order of events AND their timing"),
    (1, "not the fastest time — the right time"),
    (0, "A controller (discrete) drives a plant (continuous)"),
    (0, "Everyday examples"),
    (1, "airbags, cruise control, production lines, robots, protocols"),
], gap=8)
callout(s, 8.15, 1.7, 4.55, 2.7, "Why proof, not vibes",
        ["These are safety-critical: “probably correct” is not acceptable when an airbag must fire within milliseconds.",
         "Timing correctness must be proven, and that is exactly what real-time model checking does."])
pagenum(s, 12)

# 13 — Real-time model checking + workflow (merge 14+15)
s = slide()
title(s, "Real-time model checking: the working loop", kicker="Real-time systems")
# workflow chips
labels = ["Model", "Simulate", "Verify", "Interpret", "Refine"]
x = 0.98; y = 1.7; wchip = 2.0; gapx = 0.42
chips = []
for i, l in enumerate(labels):
    cxp = x + i * (wchip + gapx)
    fill = CARD if l != "Verify" else RGBColor(0xD9,0xEC,0xEF)
    rect(s, cxp, y, wchip, 0.7, fill=fill, line=TEAL, lw=1.1, rounded=True)
    r = rect(s, cxp, y, wchip, 0.7, fill=None, line=None)  # placeholder no-op
    _, tf = box(s, cxp, y + 0.16, wchip, 0.4)
    para(tf, l, size=13.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, after=0)
    if i < len(labels) - 1:
        cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(cxp + wchip), Inches(y + 0.35),
                                    Inches(cxp + wchip + gapx), Inches(y + 0.35))
        cn.line.color.rgb = NAVY; cn.line.width = Pt(1.4); _arrow(cn)
bullets(s, 0.98, 2.75, 7.1, [
    (0, "Model the tasks and the environment as automata, then verify the resulting network"),
    (0, "A counterexample sends you back into the simulator to see why"),
], gap=10)
callout(s, 8.2, 2.75, 4.5, 3.0, "Reframed for the AI era",
        ["The MODEL step is exactly what AI can now assist.",
         "",
         "The durable human skill is the rest of the loop:",
         "read the trace, interpret it, validate against intent,",
         "and decide what to refine. Don't outsource the",
         "judgement — only the drawing."])
pagenum(s, 13)

# 14 — Industrial scale (reuse ABB image)
s = slide()
title(s, "…at industrial scale", kicker="Real-time systems")
add_image(s, os.path.join(ASSETS, "abb-model.png"), 0.98, 2.25, w=4.5, frame=True)
bullets(s, 5.75, 1.7, 6.95, [
    (0, "A real protocol model (ABB): every box is an automaton"),
    (0, "Engineers ask questions no one can eyeball"),
    (1, "is this error state reachable?"),
    (1, "max response time between two states?"),
    (1, "is this variable always < 64?"),
    (1, "is the system guaranteed to reach this state?"),
    (0, "At this scale an AI-generated model MUST be checked — you cannot verify it by reading it"),
], gap=7, base=16.5)
pagenum(s, 14)

# 15 — Example: Light control
s = slide()
title(s, "Example: light control", kicker="Timed automata")
bullets(s, 0.98, 2.3, 11.0, [
    (0, "Wanted behaviour"),
    (1, "pressed once  →  light on"),
    (1, "pressed twice quickly  →  light gets brighter"),
    (1, "pressed again  →  light off"),
    (0, "The word “quickly” is about timing — hold that thought"),
], gap=18, base=24)
pagenum(s, 15)

# 16 — Finite state automata (redraw)
s = slide()
title(s, "Finite-state automata", kicker="Timed automata")
bullets(s, 0.98, 1.85, 6.6, [
    (0, "A finite state graph with"),
    (1, "a set of nodes (states / locations)"),
    (1, "a set of edges (transitions)"),
    (1, "a set of labels (actions)"),
], gap=14, base=21)
off = loc(s, 9.0, 3.6, "Off", initial=True)
on  = loc(s, 11.6, 3.6, "On")
edge(s, off, on, [('a', 'a')], off=-0.26, side='above')
edge(s, on, off, [('a', 'b')], off=-0.26, side='above')
diagram_caption(s, 8.2, 4.7, 4.4, "A two-state light: a turns it on, b turns it off")
pagenum(s, 16)

# 17 — FSA with variables (redraw)
s = slide()
title(s, "Finite-state automata with variables", kicker="Timed automata")
bullets(s, 0.98, 1.55, 5.9, [
    (0, "Extend FSA with variables"),
    (1, "guards — conditions on variables to enable a transition"),
    (1, "assignments — updates / resets on a transition"),
    (0, "Semantics: transition systems"),
])
off = loc(s, 8.7, 3.6, "Off", initial=True)
on  = loc(s, 11.5, 3.6, "On")
edge(s, off, on, [('g', 'Flag==1'), ('a', 'a'), ('u', 'Flag:=2')], off=-0.30, side='above')
edge(s, on, off, [('g', 'Flag==2'), ('a', 'b'), ('u', 'Flag:=1')], off=-0.30, side='below')
diagram_caption(s, 7.4, 5.05, 5.2, "Guards gate edges; assignments update state")
legend(s)
pagenum(s, 17)

# 18 — Light control without time fails
s = slide()
title(s, "Light control without time — it doesn't work", kicker="Timed automata")
off = loc(s, 2.4, 3.4, "Off", initial=True)
lig = loc(s, 5.6, 3.4, "Light")
bri = loc(s, 8.8, 3.4, "Bright")
edge(s, off, lig)
edge(s, lig, bri)
elbow_edge(s, lig, off, drop=0.85, end_dx=-0.12)
elbow_edge(s, bri, off, drop=1.5, end_dx=0.12)
diagram_caption(s, 1.5, 5.6, 8.0, "Every transition is a single button press")
callout(s, 10.35, 2.5, 2.5, 2.5, "The gap",
        ["Nothing here distinguishes a quick second press from a slow one.",
         "We need to measure elapsed time → add a clock."], accent=TEAL, bg=CARD)
pagenum(s, 18)

# 19 — Timed automata (redraw) s22
s = slide()
title(s, "Timed automata  (Alur & Dill, 1990)", kicker="Timed automata")
bullets(s, 0.98, 1.6, 5.7, bottom=4.25, items=[
    (0, "Guard — a timing constraint, e.g.  x > 10"),
    (0, "Action — a synchronisation label, e.g.  a"),
    (0, "Clock reset — set a clock to 0, e.g.  x := 0"),
])
n1 = loc(s, 8.8, 3.5, "1", initial=True)
n2 = loc(s, 11.6, 3.5, "2")
edge(s, n1, n2, [('g', 'x>10'), ('a', 'a'), ('u', 'x:=0')], side='above')
legend(s, 8.0, 4.4)
callout(s, 0.98, 4.4, 6.0, 1.7, "AI doesn't change this",
        ["The definition of a timed automaton is mathematics — it is",
         "exactly the same in 2025 as in 1990. What AI changes is who",
         "draws the automaton, not what it means."])
pagenum(s, 19)

# 20 — Light control with a clock (redraw) s23/24
s = slide()
title(s, "Light control, with a clock", kicker="Timed automata")
off = loc(s, 2.6, 3.4, "Off", initial=True)
lig = loc(s, 6.0, 3.4, "Light")
bri = loc(s, 9.4, 3.4, "Bright")
edge(s, off, lig, [('u', 'x:=0')], side='above')
edge(s, lig, bri, [('g', 'x<=3')], side='above')
elbow_edge(s, lig, off, [('g', 'x>3')], drop=0.95, end_dx=-0.12, lab_side='above')
elbow_edge(s, bri, off, drop=1.6, end_dx=0.12)
diagram_caption(s, 1.3, 5.75, 9.4, "Presses are the events; clock x separates a quick press (x≤3 → Bright) from a slow one (x>3 → Off)")
legend(s)
callout(s, 10.9, 2.5, 2.1, 2.2, "Analysis",
        ["Can the light", "eventually become", "bright?  Ask the", "model checker."],
        accent=TEAL, bg=CARD)
pagenum(s, 20)

# 21 — TA semantics s25
s = slide()
title(s, "Timed automata: semantics", kicker="Timed automata")
n = loc(s, 2.6, 3.3, "n", initial=True)
m = loc(s, 5.7, 3.3, "m")
edge(s, n, m, [('g', 'x<=5 & y>3'), ('a', 'a'), ('u', 'x:=0')], side='above')
legend(s, 1.0, 3.95)
bullets(s, 7.0, 1.55, 5.7, [
    (0, "State = ( location, clock values )"),
    (1, "clocks x, y take real values"),
    (0, "Two kinds of transition"),
    (1, "discrete — follow an edge (instantaneous)"),
    (1, "delay — stay put; all clocks advance by the same real amount"),
], gap=7, base=16.5)
_, tf = box(s, 0.98, 4.55, 7.6, 1.7)
para(tf, "( n, x=2.4, y=3.14 ) —a→ ( m, x=0, y=3.14 )   discrete", size=13.5, color=INK, font=MONO, first=True, after=6)
para(tf, "( n, x=2.4, y=3.14 ) —1.1→ ( n, x=3.5, y=4.24 )   delay", size=13.5, color=INK, font=MONO, after=0)
pagenum(s, 21)

# 22 — Invariants s26
s = slide()
title(s, "Invariants force progress", kicker="Timed automata")
n = loc(s, 3.0, 3.5, "n", initial=True, tag="x<=5")
m = loc(s, 6.4, 3.5, "m", tag="y<=10")
edge(s, n, m, [('g', 'x<=5 & y>3'), ('a', 'a'), ('u', 'x:=0')], side='below')
bullets(s, 7.4, 1.7, 5.35, [
    (0, "An invariant bounds how long you may stay in a location"),
    (1, "upper bounds only, e.g. x<=5"),
    (0, "It cannot delay past x=5, so it must act"),
    (1, "risk: it can also get stuck (deadlock) — a design smell to check"),
], gap=9, base=16.5)
diagram_caption(s, 1.5, 5.35, 6.5, "Without the invariant the automaton could delay forever in n")
legend(s)
pagenum(s, 22)

# 23 — Clock constraints s27
s = slide()
title(s, "Clock constraints", kicker="Timed automata")
bullets(s, 0.98, 2.2, 11.5, [
    (0, "For clocks x, y and an integer constant c, a constraint is built from:"),
    (1, "comparisons:  x < c,  x <= c,  x >= c,  x > c"),
    (1, "clock differences:  x - y <= c"),
    (1, "conjunction:  g1 & g2   (and, via negation, disjunction)"),
    (0, "Guards use these freely; invariants use upper bounds only"),
], gap=18, base=21)
pagenum(s, 23)

# 24 — TA examples: guards & invariants (merge 28-33)
s = slide()
title(s, "Timed automata: guards & invariants", kicker="Timed automata")
# three single-location self-loop examples
def loop_example(cx, title_txt, guard, tag=None, cap=None):
    node = loc(s, cx, 4.05, "l", initial=True, tag=tag, tag_side='right')
    self_loop_curved(s, node, [('g', guard), ('u', 'x:=0')])
    _, tf = box(s, cx - 1.5, 1.75, 3.0, 0.4)
    para(tf, title_txt, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, after=0)
    if cap:
        diagram_caption(s, cx - 1.6, 4.95, 3.2, cap)
loop_example(2.8, "Lower bound", "x>=2", cap="fires once x≥2, may wait longer")
loop_example(6.7, "Bounded", "2<=x<=3", cap="fires only while 2≤x≤3")
loop_example(10.6, "With invariant", "x>=2", tag="x<=3", cap="invariant x≤3 forces the action")
legend(s)
pagenum(s, 24)

# 25 — Task models (merge 34-36)
s = slide()
title(s, "Modelling tasks: periodic, sporadic, aperiodic", kicker="Timed automata")
def task(cx, ttl, guard, tag, cap):
    node = loc(s, cx, 4.05, "T", initial=True, tag=tag, tag_side='right')
    self_loop_curved(s, node, [('g', guard), ('u', 'x:=0')])
    _, tf = box(s, cx - 1.6, 1.75, 3.2, 0.4)
    para(tf, ttl, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, after=0)
    diagram_caption(s, cx - 1.7, 4.95, 3.4, cap)
task(2.8, "Periodic (20)", "x==20", "x<=20", "exactly every 20 time units")
task(6.7, "Sporadic (≥20)", "x>=20", None, "at least 20 apart, may wait")
task(10.6, "Aperiodic (5–100)", "5<=x<=99", "x<=100", "between 5 and 100 apart")
legend(s)
pagenum(s, 25)

# 26 — Light switch s37
s = slide()
title(s, "Timed automata: a light switch", kicker="Timed automata")
off = loc(s, 3.4, 4.0, "off", initial=True)
on  = loc(s, 7.4, 4.0, "on", tag="x<=9", tag_side='right')
edge(s, off, on, [('g', 'x>2'), ('a', 'push'), ('u', 'x:=0')], side='above')
elbow_edge(s, on, off, [('g', 'x==9'), ('u', 'x:=0')], drop=0.95, lab_side='below')
self_loop_curved(s, on, [('a', 'push'), ('u', 'x:=0')])
bullets(s, 9.7, 2.2, 3.0, bottom=5.25, items=[
    (0, "On only after >2 units since off"),
    (0, "Auto-off after 9 units, unless pressed again"),
], gap=8, base=16.5)
legend(s, 9.7, 5.5)
pagenum(s, 26)

# 27 — Semantics definition + trace (39/40)
s = slide()
title(s, "Semantics: action & delay transitions", kicker="Timed automata")
bullets(s, 0.98, 1.6, 6.4, bottom=4.15, items=[
    (0, "Action transition  (l, v) —a→ (l', v')"),
    (1, "guard g holds in v; apply reset r; invariant of l' holds in v'"),
    (0, "Delay transition  (l, v) —d→ (l, v+d)"),
    (1, "invariant of l must hold throughout the delay"),
], gap=8, base=17)
_, tf = box(s, 0.98, 4.3, 8.0, 1.9)
para(tf, "Example trace (light switch, clocks x, y):", size=12.5, color=MUTED, italic=True, first=True, after=5)
for ln in ["(off, 0, 0) —3.5→ (off, 3.5, 3.5)",
           "        —push→ (on, 0, 0)   [both reset]",
           "        —π→ (on, π, π)  —push→ (on, 0, π)  …"]:
    para(tf, ln, size=13, color=INK, font=MONO, after=2)
pagenum(s, 27)

# 28 — Networks of timed automata s41
s = slide()
title(s, "Networks of timed automata", kicker="Timed automata")
bullets(s, 0.98, 1.55, 5.4, bottom=5.95, items=[
    (0, "Several automata run in parallel"),
    (0, "They synchronise on channels"),
    (1, "a!  — send   /   a?  — receive"),
    (1, "both edges are taken together, atomically"),
    (0, "Shared (bounded) integer variables too"),
])
a1 = loc(s, 8.4, 2.7, "l1", initial=True)
a2 = loc(s, 11.4, 2.7, "l2")
b1 = loc(s, 8.4, 5.1, "m1", initial=True)
b2 = loc(s, 11.4, 5.1, "m2")
edge(s, a1, a2, [('g', 'x>=2 & i==3'), ('a', 'a!'), ('u', 'y:=0')], side='above')
edge(s, b1, b2, [('a', 'a?'), ('u', 'i:=i+1')], side='below')
diagram_caption(s, 7.4, 5.95, 5.2, "a! and a? fire together — a two-way rendezvous")
legend(s, 0.98, 6.2)
pagenum(s, 28)

# 29 — How to specify what to check s43
s = slide()
title(s, "Specifying requirements", kicker="Specification")
add_image(s, os.path.join(ASSETS, "tctl-tree.png"), 8.0, 1.7, w=4.6)
bullets(s, 0.98, 1.7, 6.6, [
    (0, "Requirements become temporal-logic formulae"),
    (1, "Invariant — something bad never happens"),
    (1, "Reachability — something may happen"),
    (1, "Liveness — something eventually happens"),
    (0, "Choosing the RIGHT property is validation"),
    (1, "AI can draft a formula; only you know which requirement it must capture"),
], gap=8, base=17)
pagenum(s, 29)

# 30 — TCTL s44
s = slide()
title(s, "TCTL — Timed Computation Tree Logic", kicker="Specification")
bullets(s, 0.98, 1.7, 7.0, [
    (0, "A branching-time logic over the model's computation tree"),
    (1, "state formulae describe a single state"),
    (1, "path formulae quantify over paths (traces)"),
    (0, "A path carrying timing is a timed path"),
    (0, "Path formulae: reachability, safety, liveness"),
], gap=9, base=17)
callout(s, 8.5, 1.9, 4.2, 2.1, "UPPAAL",
        ["UPPAAL verifies a decidable subset of TCTL — expressive enough for the real-time properties engineers actually need to check."],
        accent=TEAL, bg=CARD)
pagenum(s, 30)

# 31 — Quantifiers s45
s = slide()
title(s, "Quantifiers in TCTL", kicker="Specification")
bullets(s, 0.98, 1.7, 6.2, [
    (0, "Path quantifiers"),
    (1, "E — there exists a path"),
    (1, "A — for all paths"),
    (0, "State quantifiers along a path"),
    (1, "[]  — all states  (G)"),
    (1, "<> — some state  (F)"),
    (0, "Combinations we use:  A[], A<>, E<>, E[]"),
], gap=7, base=17)
callout(s, 7.9, 1.9, 4.8, 2.4, "In words",
        ["A[] p   — invariantly p   (always)",
         "E<> p   — p is reachable   (possibly)",
         "A<> p   — p is inevitable  (unavoidable)",
         "E[] p   — potentially always p"],
        accent=TEAL, bg=CARD)
pagenum(s, 31)

# 32 — Property patterns (merge 46-50)
s = slide()
title(s, "The property patterns you'll use", kicker="Specification")
rows = [
    ("E<> p", "Reachable", "p holds in at least one reachable state — sanity checks"),
    ("A[] p", "Invariant / safety", "p holds in all reachable states — “bad never happens”"),
    ("E[] p", "Potentially always", "some path keeps p true forever"),
    ("A<> p", "Inevitable / liveness", "every path eventually reaches a p-state"),
    ("p --> q", "Leads-to", "A[](p imply A<> q) — whenever p, eventually q"),
]
y = 1.75
for form, name, desc in rows:
    rect(s, 0.98, y, 2.0, 0.72, fill=CARD, line=TEAL, lw=1.0, rounded=True)
    _, tf = box(s, 1.02, y + 0.2, 1.92, 0.4)
    para(tf, form, size=14, color=NAVY, bold=True, font=MONO, align=PP_ALIGN.CENTER, first=True, after=0)
    _, tf = box(s, 3.2, y + 0.06, 3.0, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, size=13.5, color=NAVY, bold=True, first=True, after=0)
    _, tf = box(s, 6.3, y + 0.06, 6.4, 0.66, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, desc, size=12.5, color=MUTED, first=True, after=0, line=1.02)
    y += 0.86
pagenum(s, 32)

# 33 — Using AI to build models (NEW block 1)
s = slide()
title(s, "Using AI to build formal models", kicker="AI for formal modelling")
colw = 5.7
rect(s, 0.98, 1.7, colw, 3.7, fill=RGBColor(0xE8,0xF3,0xEC), line=RGBColor(0x2E,0x8B,0x57), lw=1.2, rounded=True)
_, tf = box(s, 1.25, 1.9, colw-0.5, 3.4)
para(tf, "AI is good at", size=15, color=RGBColor(0x2E,0x8B,0x57), bold=True, first=True, after=6)
for ln in ["drafting a model from a description",
           "boilerplate and UPPAAL syntax",
           "candidate TCTL properties",
           "quick iterations on your feedback"]:
    para(tf, "•  " + ln, size=13.5, color=INK, after=5)
rect(s, 7.05, 1.7, colw, 3.7, fill=AMBERBG, line=AMBER, lw=1.2, rounded=True)
_, tf = box(s, 7.32, 1.9, colw-0.5, 3.4)
para(tf, "AI is bad at", size=15, color=AMBER, bold=True, first=True, after=6)
for ln in ["guarantees — it gives plausible, not proven",
           "hallucinated counterexample traces",
           "missing subtle timing edge cases",
           "the intent you didn't state (validation)"]:
    para(tf, "•  " + ln, size=13.5, color=INK, after=5)
_, tf = box(s, 0.98, 5.6, 11.7, 0.7)
para(tf, "Maps straight onto verification (AI helps) vs validation (yours). AI manipulates the symbols; you supply the intent and demand the guarantee.",
     size=13.5, color=NAVY, italic=True, first=True, after=0, line=1.05)
pagenum(s, 33)

# 34 — How to prompt for a UPPAAL model (NEW block 2)
s = slide()
title(s, "How to prompt for a UPPAAL model", kicker="AI for formal modelling")
bullets(s, 0.98, 1.6, 7.1, [
    (0, "State the requirement precisely — units, bounds, edge cases"),
    (0, "Name the formalism: a network of timed automata"),
    (0, "Enumerate the components and their shared events (channels)"),
    (0, "Ask for a single runnable .xml with the queries embedded"),
    (0, "Iterate on the counterexamples the tool returns"),
], gap=9, base=17)
callout(s, 8.35, 1.75, 4.35, 3.6, "Point your AI at a UPPAAL skill",
        ["A UPPAAL skill (one you build — next slide) encodes the modelling language and the exact .xml file format, so the AI emits models that actually load and verify rather than plausible-looking guesswork.",
         "Good prompt + skill → a correct model. You still own the question: is it the RIGHT model?"], accent=TEAL, bg=CARD)
pagenum(s, 34)

# 35 — NEW: Build your own skill
s = slide()
title(s, "Build your own skill", kicker="AI for formal modelling")
bullets(s, 0.98, 1.6, 6.1, bottom=5.35, items=[
    (0, "A bare LLM drifts on UPPAAL's exact syntax and file format"),
    (1, "capture the know-how in a skill → models that load and verify"),
    (0, "An Agent Skill = a small folder your AI loads on demand"),
    (1, "a SKILL.md: when to use it, and the procedure"),
    (1, "reference docs it opens only when needed"),
    (1, "portable across AI tools (Claude Code, Codex, …)"),
], gap=7, base=16.5)
# right card: the four build steps
rect(s, 7.15, 1.6, 5.55, 3.75, fill=CARD, line=TEAL, lw=1.2, rounded=True)
rect(s, 7.15, 1.6, 0.11, 3.75, fill=TEAL)
_, tf = box(s, 7.42, 1.8, 5.15, 3.4)
para(tf, "BUILD A UPPAAL SKILL", size=12, color=TEAL, bold=True, first=True, after=7)
for n, ln in enumerate([
    "Ground it in the official docs (docs.uppaal.org): the modelling language and the .xml format.",
    "Write a thin SKILL.md: the modelling workflow + correctness rules (escaping, upper-bound invariants, embedded queries).",
    "Add a reference doc: the concrete .xml structure and a small worked example.",
    "Validate: generate a model, load it in UPPAAL, confirm it verifies.",
], start=1):
    p = para(tf, "", size=12.5, color=INK, first=False, after=6, line=1.06)
    r = p.add_run(); r.text = f"{n}.  "; _set(r, 12.5, TEAL, True)
    r2 = p.add_run(); r2.text = ln; _set(r2, 12.5, INK)
callout(s, 0.98, 5.55, 11.72, 1.2, "We don't hand you the skill — building it is the point",
        ["To make an AI produce correct UPPAAL models you must understand UPPAAL yourself — and be able to check its output.",
         "That understanding is exactly what the assignment rewards."])
pagenum(s, 35)

# 36 — Worked example: bridge problem (reuse image + the actual prompt)
s = slide()
title(s, "Worked example: the bridge problem", kicker="AI for formal modelling")
add_image(s, os.path.join(ASSETS, "bridge.png"), 0.98, 1.6, w=4.6)
# the natural-language prompt, shown as a prompt card
pcard = rect(s, 0.98, 3.75, 5.75, 2.95, fill=CARD, line=TEAL, lw=1.2, rounded=True)
rect(s, 0.98, 3.75, 0.11, 2.95, fill=TEAL)
_, tf = box(s, 1.24, 3.9, 5.35, 2.7)
para(tf, "A GOOD PROMPT — PLAIN ENGLISH + THE QUESTION", size=11, color=TEAL, bold=True, first=True, after=6)
for ln in [
    "“Using your UPPAAL skill, build a model for the",
    "bridge-crossing puzzle. Four people cross a narrow",
    "bridge at night with one torch; it holds two at a time,",
    "so someone must carry the torch back; a pair moves at",
    "the slower one's pace. Crossing times: 5, 10, 20, 25 min.",
    "Can all four cross within 60 minutes? Add a query that",
    "answers this, verify it, and check it cannot deadlock.”",
]:
    para(tf, ln, size=10.5, color=INK, font=MONO, after=1, line=1.08)
bullets(s, 7.05, 1.6, 5.6, [
    (0, "The prompt gives the problem and the question — not the automaton"),
    (0, "Your UPPAAL skill supplies the modelling"),
    (1, "locations, a clock, the guards, and the TCTL query"),
    (0, "UPPAAL then confirms a schedule exists"),
    (1, "`E<> (all safe && total <= 60)`"),
    (0, "AI built it fast; the model checker made it trustworthy"),
], gap=8, base=16.5)
pagenum(s, 36)

# 36 — Train crossing (reuse image) s54
s = slide()
title(s, "Another classic: the train crossing", kicker="Worked examples")
add_image(s, os.path.join(ASSETS, "train-crossing.png"), 0.98, 1.75, w=5.4)
bullets(s, 6.75, 1.8, 5.9, [
    (0, "Several trains, one shared bridge — mutual exclusion under timing"),
    (1, "approach signal, a 10-unit window to stop, then the crossing"),
    (1, "a controller issues stop / go as trains leave"),
    (0, "Exactly the kind of concurrent, timed coordination where a hand argument is hopeless and model checking shines"),
], gap=9, base=16.5)
pagenum(s, 37)

# 37 — UPPAAL SMC (merge 56+57)
s = slide()
title(s, "When behaviour is stochastic: UPPAAL SMC", kicker="Worked examples")
bullets(s, 0.98, 1.7, 7.1, [
    (0, "Statistical model checking for stochastic timed automata"),
    (1, "probabilistic choice between enabled transitions"),
    (1, "distributions over delays; ODEs for continuous variables"),
    (0, "Queries answer probabilistic questions"),
    (1, "estimate Pr[<=T](<> p);  test Pr >= p0;  compare two probabilities"),
], gap=8, base=17)
callout(s, 8.4, 1.85, 4.3, 2.0, "Use when the requirement is probabilistic",
        ["Classic TCTL proves deterministic timing;",
         "SMC estimates likelihoods when the system itself",
         "is random. Different question, same discipline."], accent=TEAL, bg=CARD)
pagenum(s, 38)

# 38 — Reflection: what's yours vs the AI's (NEW, dark)
s = slide(dark=True)
title(s, "What's yours vs. what's the AI's", dark=True, kicker="Before the assignment")
rect(s, 0.95, 2.1, 5.7, 3.9, fill=RGBColor(0x22,0x2A,0x52), line=RGBColor(0x3E,0x9A,0xA8), lw=1.2, rounded=True)
_, tf = box(s, 1.25, 2.35, 5.1, 3.5)
para(tf, "The AI can", size=16, color=RGBColor(0x86,0xD6,0xC9), bold=True, font=HEAD, first=True, after=8)
for ln in ["draft the automata", "write UPPAAL syntax", "suggest properties", "run the tool"]:
    para(tf, "•  " + ln, size=14, color=RGBColor(0xC7,0xD2,0xE8), after=5)
rect(s, 6.9, 2.1, 5.5, 3.9, fill=RGBColor(0x2A,0x24,0x18), line=AMBER, lw=1.2, rounded=True)
_, tf = box(s, 7.2, 2.35, 4.9, 3.5)
para(tf, "Only you can", size=16, color=RGBColor(0xE7,0xB6,0x6A), bold=True, font=HEAD, first=True, after=8)
for ln in ["decide what to model (intent)",
           "choose the right property",
           "validate the model against reality",
           "judge whether the counterexample means the model or the requirement is wrong"]:
    para(tf, "•  " + ln, size=14, color=RGBColor(0xDD,0xD3,0xC2), after=5)
pagenum(s, 39, dark=True)

# 39 — Assignment: scenario
s = slide()
title(s, "Assignment: a pedestrian crossing controller", kicker="Assignment")
bullets(s, 0.98, 1.55, 11.6, bottom=5.15, items=[
    (0, "Model a signalised pedestrian crossing: vehicle lights, WALK / DON'T-WALK, a push button"),
    (0, "The requirement is deliberately under-specified — you must pin down and justify the numbers"),
    (1, "pedestrians “shouldn't wait too long”; WALK must last “long enough” to cross"),
    (1, "cars and pedestrians must NEVER get a green at the same time"),
    (1, "respond promptly to a press, but don't switch too rapidly"),
    (0, "You receive an assigned parameter set / twist (crossing width, a bus-priority input, a night mode)"),
], gap=12, base=17)
callout(s, 0.98, 5.35, 11.7, 1.25, "Why this is AI-resistant",
        ["The thresholds are yours to choose and defend against real-world intent — the AI never received that intent, so it",
         "cannot justify them for you. Paste the brief into any tool and the judgement, traces and defence remain your work."])
pagenum(s, 40)

# 40 — Assignment: deliverables
s = slide()
title(s, "Assignment: what you submit (graded)", kicker="Assignment")
bullets(s, 0.98, 1.5, 11.7, [
    (0, "Requirement interpretation — resolve each ambiguity; state your thresholds and justify them against intent"),
    (0, "Formalisation — the requirements as TCTL properties, with a reason for each interpretation"),
    (0, "Model + emergent-failure hunt — build it (AI-assisted is fine); use verification to find where it breaks"),
    (1, "document each failure with the counterexample trace UPPAAL produces"),
    (0, "Hand-traced counterexample — reproduce one failing trace by hand: locations and clock valuations, step by step"),
    (0, "Fix & re-prove — correct the model and show the property now holds"),
    (0, "Reflection — where AI helped, where it was wrong, and how verification/validation caught it"),
], gap=7, base=16.5)
pagenum(s, 41)

# 41 — Assignment: submission + AI stance
s = slide()
title(s, "Assignment: submission & using AI", kicker="Assignment")
bullets(s, 0.98, 1.7, 7.1, [
    (0, "Submit"),
    (1, "one runnable UPPAAL model .xml with queries embedded"),
    (1, "the counterexample traces and your hand trace"),
    (1, "a short report covering all seven deliverables"),
    (0, "Using AI is allowed and expected"),
    (1, "the grade is what you add on top of its output"),
], gap=8, base=17)
callout(s, 8.35, 1.85, 4.35, 3.4, "The test we designed against",
        ["“If you paste the whole brief into any AI tool, what is left for you to do?”",
         "Answer: choose and defend the thresholds, hand-trace a counterexample, and validate the model against the real crossing.",
         "That is the part that shows you learned — and the AI cannot do it for you."])
pagenum(s, 42)

# ---------- speaker notes (delivery script, new verification/validation flow) ----------
NOTES = [
"Welcome. This is Formal Modeling and Analysis of Real-Time Systems. We teach it this year with an explicit question running through it: now that AI can build models for us, what is actually worth learning? Short answer — the concepts, and the judgement to check what the AI produces.",
"A quick word on where I'm coming from: embedded software, a PhD and postdoc on formal modelling and analysis of embedded systems, and now formal methods at Prover.",
"Start with the basic idea of a model: an abstraction of reality that helps us understand it. The hook on the right is the theme of the course — an LLM is itself a model of the world, but a statistical, non-semantic one. You can't hand it to a prover and get a guarantee. We build models we CAN reason about.",
"Three perspectives — external, behavioural, structural — and two classic example families. Good discussion prompt: 'is an AI a model?' Yes, but not a verifiable one. That difference is the whole point.",
"Formal = mathematical semantics: every symbol has a precise meaning. Emphasise the trust anchor framing — semantics is exactly what lets you check someone else's model, including one an AI generated. No semantics, no proof, no basis for trust.",
"This is the spine of the course. Verification asks 'did we build the thing right?' — symbolic, exhaustive, automatable, and where AI helps most. Validation asks 'did we build the right thing?' — it connects the model to real-world intent and lives OUTSIDE the formal system, so AI can't own it. Repeat the line: AI manipulates the symbols; it cannot supply the intent or the guarantee.",
"Formal verification checks ALL behaviours, exhaustively, against a property in logic. Contrast with an LLM, which gives a plausible answer fast but with no guarantee. Powerful is not the same as certain.",
"The catch is state-space explosion — the number of states blows up with components, clocks and variables. Make the point clearly: this is an open problem that has nothing to do with how good LLMs get. It's the concrete rebuttal to 'why model-check at all now?'",
"The thesis slide. Walk the three cards: AI lowers the historic cost of formal methods; that raises the stakes, because now you must verify AI's output and AI is a fallible author; so your job shifts from drawing automata to specifying, validating, and judging counterexamples. That last skill is the durable one.",
"Deductive verification proves correctness from axioms and rules — not fully automated. The Socrates syllogism is the toy example. Honest note: this is genuinely shifted by AI — LLMs plus proof assistants are pushing autoformalization. But imitating reasoning is not the same as a sound proof.",
"Model checking takes two inputs: the model A (a network of timed automata) and the requirement F (a temporal-logic formula). The checker is the oracle — it answers with certainty and hands you a counterexample when F fails. It is where an AI-built model gets checked.",
"Real-time systems: correctness depends on order AND timing — the right time, not the fastest. A discrete controller drives a continuous plant. These are safety-critical, so 'probably correct' is not acceptable; timing must be proven.",
"The working loop: model, simulate, verify, interpret, refine. Here's the reframing — the MODEL step is exactly what AI now assists. The durable human skill is the rest: read the trace, interpret it, validate against intent, decide what to fix. Outsource the drawing, not the judgement.",
"A real industrial protocol model from ABB — every box is an automaton. Engineers ask questions no one can answer by eye: is this error state reachable, what's the max response time, is this variable always under 64. At this scale an AI-generated model MUST be checked; you cannot verify it by reading it.",
"Our running example. The wanted behaviour is simple to state — but note the word 'quickly'. That's a timing requirement, and it's what forces us toward clocks.",
"Finite-state automata: nodes (states), edges (transitions), labels (actions). The minimal foundation.",
"Extend FSA with variables: guards decide when an edge is enabled, assignments update state on the edge. Semantics are transition systems.",
"Try to model light control with a plain FSA and it fails — nothing distinguishes a quick second press from a slow one. Every edge is just 'press'. We need to measure elapsed time, so we add a clock.",
"Timed automata, Alur and Dill 1990: guards are timing constraints, actions synchronise, clocks reset to zero. Say it plainly — this definition is mathematics; it's the same in 2025 as in 1990. AI changes who draws the automaton, not what it means.",
"Now light control works: clock x measures the delay between presses. A quick press (x<=3) goes to Bright; a slow one (x>3) returns to Off. The natural analysis question — can the light eventually become bright — is one we hand to the model checker.",
"Semantics: a state is a location plus real-valued clock readings. Two transition kinds — discrete (follow an edge, instantaneous) and delay (stay put, all clocks advance by the same real amount). The two example lines show each.",
"Invariants are upper bounds that force progress: the automaton can't delay past the bound, so it must act. Flag the risk — invariants can also cause deadlock, which is itself something to verify.",
"The grammar of clock constraints: comparisons, clock differences, conjunction. Guards use them freely; invariants use upper bounds only.",
"Three variants on one location: a lower-bound guard, a bounded window, and a guard plus an invariant that forces the action. This is where students build intuition for hand-tracing.",
"Modelling task arrivals: periodic fires exactly every 20 (guard x==20, invariant x<=20); sporadic is at-least-20-apart (guard x>=20, no invariant); aperiodic is between 5 and 100 (guard 5<=x<=99, invariant x<=100).",
"A fuller example, the light switch: on only after more than 2 units since off, auto-off after 9, and a re-press resets the clock (the self-loop). Good one to pause on.",
"The formal rules behind it: an action transition needs its guard satisfied and the target invariant to hold after the reset; a delay transition keeps the invariant true throughout. The trace at the bottom is exactly the kind of hand-tracing the assignment will ask for.",
"Compose several automata into a network; they synchronise on channels — a! sends, a? receives, and both edges fire together atomically. Shared bounded integers too. This is how we model interacting components.",
"Now: how do we say what to check? Requirements become temporal-logic formulae — invariant, reachability, liveness. Stress that choosing the RIGHT property is validation: AI can draft a formula, but only you know which requirement it has to capture.",
"TCTL is a branching-time logic over the model's computation tree: state formulae describe one state, path formulae quantify over traces, and a trace with timing is a timed path. UPPAAL checks a decidable subset.",
"The quantifiers: E/A over paths, [] and <> over states along a path, and the four combinations we actually use. The right box gives the plain-English reading.",
"The five patterns students will use in the lab. E<> for sanity/reachability, A[] for safety, A<> for liveness, E[] for potentially-always, and leads-to for response properties. Point out leads-to is just A[](p imply A<> q).",
"Now we turn to using AI for the modelling itself. Left column: what it's good at — drafting models, syntax, candidate properties, fast iteration. Right: what it's bad at — guarantees, hallucinated traces, missed edge cases, and the intent you didn't state. It maps straight onto verification versus validation.",
"The practical recipe for prompting: precise requirement with units and edge cases, name the formalism, list components and channels, ask for a single runnable .xml with queries embedded, and iterate on counterexamples. Point them at a UPPAAL skill — which they build in the next slide — so the AI emits models that actually load and verify.",
"Key twist: we do NOT give you a ready-made UPPAAL skill — if you want one, you build it, and that's the point. Explain what an Agent Skill is: a small folder your AI loads on demand, with a SKILL.md that says when to use it and how, plus reference docs it opens only when needed. Walk the four build steps on the right: ground it in the official UPPAAL docs, write a thin SKILL.md with the workflow and correctness rules, add a reference with the concrete .xml format and an example, then validate by generating a model and checking it verifies in UPPAAL. Land the message: to make an AI produce correct models you must understand UPPAAL yourself — which is exactly what the assignment rewards.",
"Worked example — the bridge problem. Read the prompt on the slide out loud, and make the point: it's plain English — it states the puzzle and the question, but says nothing about locations, clocks, or guards. That's what a good prompt looks like. The UPPAAL skill supplies the modelling, and UPPAAL confirms a 60-minute schedule exists. The AI built it fast; the model checker made it trustworthy.",
"The train crossing: several trains, one shared bridge, mutual exclusion under timing constraints and a controller. Exactly the concurrent, timed coordination where a hand argument is hopeless and model checking earns its keep.",
"When behaviour is genuinely random, UPPAAL SMC handles stochastic timed automata and answers probabilistic questions — estimate a probability, test it against a threshold, compare two. Different question, same discipline. Keep this brief.",
"Reflection before the assignment. Left: what the AI can do for you. Right: what only you can do — decide what to model, choose the right property, validate against reality, and judge whether a counterexample means the model or the requirement is wrong. That right-hand column is what we grade.",
"The assignment: a pedestrian crossing controller. It's deliberately under-specified — 'not too long', 'long enough' — so you must fix the numbers and justify them. Everyone gets an assigned parameter set or twist. Explain why that makes it AI-resistant: the AI never received the intent, so it can't defend your thresholds for you.",
"Walk through the seven graded deliverables. The load-bearing ones are the threshold justification, the emergent-failure hunt with real counterexample traces, and especially the hand-traced counterexample — locations and clock valuations, step by step — which an AI hallucinates.",
"Submission is a single runnable .xml with queries embedded. Using AI is allowed and expected — the grade is what you add on top of its output. End on the design test: if you paste the whole brief into any tool, what's left for you to do? The parts that show you learned.",
]
for _i, _sld in enumerate(prs.slides):
    if _i < len(NOTES) and NOTES[_i]:
        _sld.notes_slide.notes_text_frame.text = NOTES[_i]

# ---------- save ----------
out = os.path.join(HERE, "Slides-AI-edition.pptx")
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
